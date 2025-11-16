#!/usr/bin/env python3
"""
Universal PowerPoint reader supporting both legacy .ppt and modern .pptx formats.

Features:
- Reads .pptx files directly with python-pptx
- Converts legacy .ppt files to .pptx using LibreOffice if available
- Extracts text content, slide structure, and speaker notes
- Saves structured JSON output for downstream processing

Usage:
    python tools/read_powerpoint.py <path_to_file.ppt|pptx>
    python tools/read_powerpoint.py <path_to_file.ppt|pptx> --output custom_output.json
"""

import argparse
import json
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Optional, Any

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")


def check_file_type(file_path: Path) -> str:
    """
    Check if file is legacy .ppt or modern .pptx format.
    
    Args:
        file_path: Path to PowerPoint file
        
    Returns:
        'ppt' for legacy format, 'pptx' for modern format, 'unknown' for others
    """
    try:
        result = subprocess.run(
            ['file', str(file_path)], 
            capture_output=True, 
            text=True, 
            timeout=5
        )
        output = result.stdout.lower()
        
        if 'composite document' in output or 'ole 2 compound document' in output:
            return 'ppt'
        elif 'zip archive' in output or 'microsoft ooxml' in output:
            return 'pptx'
        else:
            return 'unknown'
    except Exception as e:
        print(f"Warning: Could not determine file type: {e}")
        return 'unknown'


def convert_ppt_to_pptx(ppt_path: Path, keep_original: bool = True) -> Optional[Path]:
    """
    Convert legacy .ppt to modern .pptx using LibreOffice.
    
    Args:
        ppt_path: Path to .ppt file
        keep_original: If True, keeps the original .ppt file
        
    Returns:
        Path to converted .pptx file, or None if conversion failed
    """
    output_dir = ppt_path.parent
    output_file = output_dir / f"{ppt_path.stem}_converted.pptx"
    
    # Check if LibreOffice is available
    libreoffice_commands = ['soffice', 'libreoffice', '/Applications/LibreOffice.app/Contents/MacOS/soffice']
    libreoffice_cmd = None
    
    for cmd in libreoffice_commands:
        try:
            result = subprocess.run([cmd, '--version'], capture_output=True, timeout=5)
            if result.returncode == 0:
                libreoffice_cmd = cmd
                break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    
    if not libreoffice_cmd:
        print("❌ LibreOffice not found. Install with: brew install --cask libreoffice")
        return None
    
    print(f"🔄 Converting {ppt_path.name} to .pptx format...")
    
    try:
        cmd = [
            libreoffice_cmd,
            '--headless',
            '--convert-to', 'pptx',
            '--outdir', str(output_dir),
            str(ppt_path)
        ]
        
        result = subprocess.run(
            cmd, 
            capture_output=True, 
            text=True, 
            timeout=120  # 2 minutes for large files
        )
        
        # LibreOffice creates file with original name, rename it
        default_output = output_dir / f"{ppt_path.stem}.pptx"
        if default_output.exists() and default_output != ppt_path:
            if output_file.exists():
                output_file.unlink()
            default_output.rename(output_file)
            print(f"✅ Converted to: {output_file.name}")
            return output_file
        else:
            print(f"❌ Conversion failed: {result.stderr}")
            return None
            
    except subprocess.TimeoutExpired:
        print("❌ Conversion timed out (file may be too large)")
        return None
    except Exception as e:
        print(f"❌ Conversion error: {e}")
        return None


def read_pptx(file_path: Path) -> Dict[str, Any]:
    """
    Read .pptx file and extract structured content.
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        Dict with slides, titles, content, notes, and metadata
    """
    if not HAS_PPTX:
        return {
            "error": "python-pptx not installed",
            "file_path": str(file_path)
        }
    
    try:
        prs = Presentation(file_path)
        
        data = {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "total_slides": len(prs.slides),
            "slides": [],
            "metadata": {
                "core_properties": {}
            }
        }
        
        # Extract core properties if available
        try:
            core_props = prs.core_properties
            data["metadata"]["core_properties"] = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
            }
        except:
            pass
        
        # Extract slides
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "content": [],
                "notes": "",
                "shapes": [],
                "tables": []
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Store in content list
                    slide_data["content"].append(text)
                    
                    # Detect title (first placeholder or title shape)
                    if not slide_data["title"]:
                        if hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type == 1:  # Title
                                slide_data["title"] = text
                        elif slide_num == 1 or "title" in shape.name.lower():
                            slide_data["title"] = text
                    
                    # Store shape info
                    slide_data["shapes"].append({
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "text": text
                    })
                
                # Extract tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_data["tables"].append(table_data)
            
            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_data["notes"] = notes_slide.notes_text_frame.text.strip()
            
            # Use first content as title if no title found
            if not slide_data["title"] and slide_data["content"]:
                slide_data["title"] = slide_data["content"][0]
            
            data["slides"].append(slide_data)
        
        return data
        
    except Exception as e:
        return {
            "error": f"Failed to read PowerPoint: {str(e)}",
            "file_path": str(file_path)
        }


def print_summary(data: Dict[str, Any]) -> None:
    """Print a formatted summary of the PowerPoint content."""
    print("\n" + "="*80)
    print("POWERPOINT SUMMARY")
    print("="*80)
    
    if "error" in data:
        print(f"❌ Error: {data['error']}")
        return
    
    print(f"File: {data.get('file_name', 'Unknown')}")
    print(f"Total Slides: {data.get('total_slides', 0)}")
    
    # Metadata
    if "metadata" in data and "core_properties" in data["metadata"]:
        props = data["metadata"]["core_properties"]
        if props.get("author"):
            print(f"Author: {props['author']}")
        if props.get("created"):
            print(f"Created: {props['created']}")
    
    print()
    
    # Slide summaries
    slides = data.get("slides", [])
    for slide in slides:
        title = slide.get("title", "Untitled")
        content_count = len(slide.get("content", []))
        notes = slide.get("notes", "")
        tables = slide.get("tables", [])
        
        print(f"Slide {slide['slide_number']}: {title[:60]}")
        print(f"  └─ {content_count} text elements, {len(tables)} tables, {'notes' if notes else 'no notes'}")


def print_detailed(data: Dict[str, Any]) -> None:
    """Print detailed content of all slides."""
    print("\n" + "="*80)
    print("DETAILED CONTENT")
    print("="*80)
    
    if "error" in data:
        print(f"❌ Error: {data['error']}")
        return
    
    for slide in data.get("slides", []):
        print(f"\n{'='*80}")
        print(f"SLIDE {slide['slide_number']}: {slide.get('title', 'Untitled')}")
        print(f"{'='*80}")
        
        # Content
        if slide.get("content"):
            print("\nContent:")
            for i, text in enumerate(slide["content"], 1):
                print(f"{i}. {text}")
        
        # Tables
        if slide.get("tables"):
            print("\nTables:")
            for i, table in enumerate(slide["tables"], 1):
                print(f"Table {i}:")
                for row in table:
                    print(f"  {' | '.join(row)}")
        
        # Notes
        if slide.get("notes"):
            print(f"\nSpeaker Notes:")
            print(slide["notes"])


def save_json(data: Dict[str, Any], output_path: Path) -> None:
    """Save extracted data to JSON file."""
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        print(f"\n✅ Saved JSON to: {output_path}")
    except Exception as e:
        print(f"\n❌ Failed to save JSON: {e}")


def main():
    parser = argparse.ArgumentParser(
        description="Read and extract content from PowerPoint files (.ppt and .pptx)"
    )
    parser.add_argument(
        "file",
        type=str,
        help="Path to PowerPoint file"
    )
    parser.add_argument(
        "-o", "--output",
        type=str,
        help="Output JSON file path (default: <filename>_content.json)"
    )
    parser.add_argument(
        "-d", "--detailed",
        action="store_true",
        help="Print detailed content of all slides"
    )
    parser.add_argument(
        "--no-convert",
        action="store_true",
        help="Skip conversion of legacy .ppt files"
    )
    
    args = parser.parse_args()
    
    # Validate input file
    file_path = Path(args.file)
    if not file_path.exists():
        print(f"❌ File not found: {file_path}")
        sys.exit(1)
    
    print(f"📄 Reading: {file_path.name}")
    
    # Check file type
    file_type = check_file_type(file_path)
    print(f"📋 Detected format: {file_type.upper()}")
    
    # Handle legacy .ppt files
    if file_type == 'ppt':
        if args.no_convert:
            print("❌ Legacy .ppt format requires conversion. Remove --no-convert flag.")
            sys.exit(1)
        
        converted_path = convert_ppt_to_pptx(file_path)
        if not converted_path:
            print("\n⚠️  Conversion failed. To convert manually:")
            print("  1. Open file in Microsoft PowerPoint")
            print("  2. File → Save As → PowerPoint Presentation (.pptx)")
            print("  3. Run this script on the converted file")
            sys.exit(1)
        
        file_path = converted_path
    
    # Read PowerPoint content
    data = read_pptx(file_path)
    
    # Print summary or detailed view
    if args.detailed:
        print_detailed(data)
    else:
        print_summary(data)
    
    # Save to JSON
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = file_path.parent / f"{file_path.stem}_content.json"
    
    save_json(data, output_path)
    
    # Final statistics
    if "slides" in data:
        print(f"\n📊 Statistics:")
        print(f"  Slides with titles: {sum(1 for s in data['slides'] if s.get('title'))}")
        print(f"  Slides with notes: {sum(1 for s in data['slides'] if s.get('notes'))}")
        print(f"  Total text blocks: {sum(len(s.get('content', [])) for s in data['slides'])}")
        print(f"  Total tables: {sum(len(s.get('tables', [])) for s in data['slides'])}")


if __name__ == "__main__":
    main()
