#!/usr/bin/env python3
"""
Generate enterprise-grade PowerPoint presentation from markdown using python-pptx.

Author: Chin-Yew Lin
Created: November 17, 2025

Requirements:
- python-pptx (pip install python-pptx)
- Python 3.8+

Usage:
    python generate_slides_pptx.py input.md output.pptx
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def parse_markdown_slides(md_file):
    """Parse markdown file into slide structure."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        # Skip metadata
        if line.startswith('**Author') or line.startswith('**Created') or line.startswith('**Format'):
            i += 1
            continue
        
        # Slide separator
        if line.startswith('---') or line.startswith('## Slide'):
            if current_slide and (current_slide['title'] or current_slide['content']):
                slides.append(current_slide)
            current_slide = {'title': '', 'content': [], 'type': 'content'}
            i += 1
            continue
        
        # Extract title
        if current_slide and not current_slide['title'] and line and not line.startswith('---'):
            current_slide['title'] = line
            i += 1
            continue
        
        # Content
        if current_slide and line:
            current_slide['content'].append(line)
        
        i += 1
    
    # Only add last slide if it has content
    if current_slide and (current_slide['title'] or current_slide['content']):
        slides.append(current_slide)
    
    return slides


def set_microsoft_colors(prs):
    """Configure Microsoft color palette."""
    # Microsoft Blue: RGB(0, 120, 212) = #0078D4
    # Microsoft Light Blue: RGB(0, 188, 242) = #00BCF2
    return {
        'primary': RGBColor(0, 120, 212),      # Microsoft Blue
        'accent': RGBColor(0, 188, 242),       # Light Blue
        'dark': RGBColor(0, 90, 158),          # Dark Blue
        'text': RGBColor(50, 49, 48),          # Dark Gray
        'light_gray': RGBColor(243, 242, 241), # Light Gray
        'green': RGBColor(16, 124, 16),        # Success Green
        'yellow': RGBColor(255, 185, 0),       # Warning Yellow
    }


def create_title_slide(prs, slide_data, colors):
    """Create title slide with gradient background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape (gradient effect using full slide width)
    bg1 = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = colors['primary']
    bg1.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0),
        Inches(11.333), Inches(2.0)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle and key message
    y_pos = 4.0
    for line in slide_data['content']:
        if '**Subtitle:**' in line:
            subtitle = line.replace('**Subtitle:**', '').strip()
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(y_pos),
                Inches(11.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            y_pos += 0.9
        
        elif '**Key Message:**' in line:
            msg = line.replace('**Key Message:**', '').strip()
            msg_box = slide.shapes.add_textbox(
                Inches(2.0), Inches(y_pos),
                Inches(9.333), Inches(1.0)
            )
            tf = msg_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = msg
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_accent_bar(slide, colors):
    """Add top accent bar to content slide."""
    bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors['primary']
    bar.line.fill.background()


def create_content_slide(prs, slide_data, slide_num, colors):
    """Create content slide with proper formatting."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add accent bar at top
    add_accent_bar(slide, colors)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.3),
        Inches(12.0), Inches(0.8)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    
    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3),
        Inches(12.0), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0)
    content_frame.margin_right = Inches(0)
    
    # Parse and add content
    is_first_paragraph = True
    for line in slide_data['content']:
        if not line.strip():
            continue
        
        # Remove markdown bold markers for plain text
        clean_line = re.sub(r'\*\*([^*]+)\*\*', r'\1', line)
        
        # Handle different line types
        if line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
            # Bullet point
            if is_first_paragraph:
                p = content_frame.paragraphs[0]
                is_first_paragraph = False
            else:
                p = content_frame.add_paragraph()
            
            p.text = clean_line[2:].strip()
            p.font.size = Pt(20)
            p.font.color.rgb = colors['text']
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            
        elif line.startswith('  - ') or line.startswith('  * '):
            # Nested bullet
            p = content_frame.add_paragraph()
            p.text = line.strip()[2:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = colors['text']
            p.level = 1
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            
        elif '**Outcome:**' in line or '**Opportunity:**' in line:
            # Special box (simplified as bold text)
            if is_first_paragraph:
                p = content_frame.paragraphs[0]
                is_first_paragraph = False
            else:
                p = content_frame.add_paragraph()
            
            p.text = clean_line
            p.font.size = Pt(19)
            p.font.bold = True
            if 'Outcome' in line:
                p.font.color.rgb = colors['green']
            else:
                p.font.color.rgb = colors['yellow']
            p.space_before = Pt(12)
            p.space_after = Pt(12)
            
        elif line.endswith(':') and not line.startswith('-'):
            # Section header
            if is_first_paragraph:
                p = content_frame.paragraphs[0]
                is_first_paragraph = False
            else:
                p = content_frame.add_paragraph()
            
            p.text = clean_line
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = colors['primary']
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            
        elif line.startswith('*(') or '(Visual:' in line:
            # Visual note
            p = content_frame.add_paragraph()
            p.text = clean_line.strip('*()')
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = RGBColor(96, 94, 92)
            p.space_before = Pt(10)
            
        else:
            # Regular paragraph
            if is_first_paragraph:
                p = content_frame.paragraphs[0]
                is_first_paragraph = False
            else:
                p = content_frame.add_paragraph()
            
            p.text = clean_line
            p.font.size = Pt(20)
            p.font.color.rgb = colors['text']
            p.space_before = Pt(8)
            p.space_after = Pt(8)
    
    # Slide number
    slide_num_box = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.0),
        Inches(0.6), Inches(0.3)
    )
    tf = slide_num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(96, 94, 92)
    p.alignment = PP_ALIGN.RIGHT
    
    return slide


def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_slides_pptx.py input.md [output.pptx]")
        print("\nExample:")
        print("  python generate_slides_pptx.py PFT_RFT_PM_Agent_Slides.md presentation.pptx")
        sys.exit(1)
    
    input_file = Path(sys.argv[1])
    if not input_file.exists():
        print(f"ERROR: Input file not found: {input_file}")
        sys.exit(1)
    
    # Output file
    if len(sys.argv) >= 3:
        output_file = Path(sys.argv[2])
    else:
        output_file = input_file.with_suffix('.pptx')
    
    print(f"Reading slides from: {input_file}")
    
    # Parse markdown
    slides = parse_markdown_slides(input_file)
    print(f"✓ Parsed {len(slides)} slides")
    
    # Create presentation with 16:9 widescreen format
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    colors = set_microsoft_colors(prs)
    
    # Add slides
    for idx, slide_data in enumerate(slides):
        if idx == 0:
            create_title_slide(prs, slide_data, colors)
            print(f"  ✓ Created title slide")
        else:
            create_content_slide(prs, slide_data, idx + 1, colors)
            print(f"  ✓ Created slide {idx + 1}: {slide_data['title'][:40]}...")
    
    # Save presentation
    prs.save(str(output_file))
    print(f"\n✓ SUCCESS! Presentation created: {output_file}")
    print(f"  Slides: {len(slides)}")
    print(f"  Size: {output_file.stat().st_size / 1024:.1f} KB")
    print(f"\nOpen with: open '{output_file}'")


if __name__ == '__main__':
    main()
