#!/usr/bin/env python3
"""
Generate PROFESSIONAL presentation with modern design using python-pptx.

Author: Chin-Yew Lin
Created: November 17, 2025

Features:
- Modern flat design with professional color scheme
- Clean typography with proper hierarchy
- Structured layouts with visual separators
- Icon-style bullet points
- Proper spacing and margins
- Professional footer with slide numbers and branding

Requirements:
- python-pptx (install: python -m pip install python-pptx)

Usage:
    python generate_slides_professional.py input.md output.pptx
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def parse_markdown_slides(md_file):
    """Parse markdown file into structured slides."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    current_slide = None
    
    for line in content.split('\n'):
        line_stripped = line.strip()
        
        # Skip metadata
        if line_stripped.startswith('**Author') or line_stripped.startswith('**Created'):
            continue
        
        # Slide separator
        if line_stripped.startswith('---') or line_stripped.startswith('## Slide'):
            if current_slide and (current_slide['title'] or current_slide['content']):
                slides.append(current_slide)
            current_slide = {'title': '', 'subtitle': '', 'content': [], 'type': 'content'}
            continue
        
        if not current_slide:
            continue
        
        # Extract title (first non-empty line after separator)
        if not current_slide['title'] and line_stripped:
            # Clean markdown formatting
            title = re.sub(r'\*\*([^*]+)\*\*', r'\1', line_stripped)
            current_slide['title'] = title
            # Detect if this is title slide
            if 'AI Project Management' in title or 'Title & Vision' in line_stripped:
                current_slide['type'] = 'title'
            continue
        
        # Extract subtitle for title slide
        if current_slide['type'] == 'title' and not current_slide['subtitle'] and line_stripped:
            subtitle = re.sub(r'\*\*([^*]+)\*\*', r'\1', line_stripped)
            if subtitle and subtitle != current_slide['title']:
                current_slide['subtitle'] = subtitle
                continue
        
        # Content
        if line_stripped:
            clean_line = re.sub(r'\*\*([^*]+)\*\*', r'\1', line_stripped)
            current_slide['content'].append(clean_line)
    
    # Add last slide
    if current_slide and (current_slide['title'] or current_slide['content']):
        slides.append(current_slide)
    
    return slides


def set_professional_colors():
    """Vibrant modern color scheme with gradients."""
    return {
        # Primary colors - More vibrant
        'navy': RGBColor(13, 27, 62),          # Deep midnight blue
        'blue': RGBColor(26, 115, 232),        # Vibrant blue (Google blue)
        'cyan': RGBColor(0, 188, 212),         # Bright cyan
        'purple': RGBColor(103, 58, 183),      # Deep purple
        'accent': RGBColor(255, 193, 7),       # Bright amber
        'orange': RGBColor(255, 87, 34),       # Vibrant orange
        
        # Text colors
        'text_dark': RGBColor(33, 33, 33),     # Almost black
        'text_gray': RGBColor(97, 97, 97),     # Medium gray
        'text_light': RGBColor(158, 158, 158), # Light gray
        
        # Background colors
        'white': RGBColor(255, 255, 255),
        'bg_gradient_1': RGBColor(26, 115, 232),  # Gradient start
        'bg_gradient_2': RGBColor(103, 58, 183),  # Gradient end
        'bg_light': RGBColor(250, 250, 250),      # Very light gray
        
        # Accent colors for highlights
        'green': RGBColor(76, 175, 80),        # Material green
        'teal': RGBColor(0, 150, 136),         # Material teal
        'pink': RGBColor(233, 30, 99),         # Material pink
    }


def add_footer(slide, slide_num, colors):
    """Add vibrant footer with gradient effect."""
    # Footer background with gradient simulation
    footer_bg1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.3),
        Inches(9), Inches(0.2)
    )
    footer_bg1.fill.solid()
    footer_bg1.fill.fore_color.rgb = colors['blue']
    footer_bg1.line.fill.background()
    
    footer_bg2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9), Inches(7.3),
        Inches(4.333), Inches(0.2)
    )
    footer_bg2.fill.solid()
    footer_bg2.fill.fore_color.rgb = colors['purple']
    footer_bg2.line.fill.background()
    
    # Slide number in circle badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(12.1), Inches(7.25),
        Inches(0.3), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = colors['accent']
    badge.line.fill.background()
    
    slide_num_box = slide.shapes.add_textbox(
        Inches(11.9), Inches(7.27),
        Inches(0.5), Inches(0.25)
    )
    text_frame = slide_num_box.text_frame
    text_frame.text = str(slide_num)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.font.name = 'Segoe UI'
    
    # Branding text (left side)
    brand_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.32),
        Inches(5), Inches(0.15)
    )
    text_frame = brand_box.text_frame
    text_frame.text = "AI Project Management Agent • PFT + RFT"
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].font.color.rgb = colors['white']
    text_frame.paragraphs[0].font.name = 'Segoe UI'


def create_title_slide(prs, slide_data, colors):
    """Create stunning title slide with gradient and geometric shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Gradient background (simulated with overlapping rectangles)
    bg1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = colors['bg_gradient_1']
    bg1.line.fill.background()
    
    # Diagonal overlay for depth
    bg2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8), Inches(0),
        Inches(5.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = colors['bg_gradient_2']
    bg2.fill.transparency = 0.3
    bg2.line.fill.background()
    bg2.rotation = 15
    
    # Large decorative circle (background element)
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1), Inches(-1),
        Inches(4), Inches(4)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = colors['cyan']
    circle1.fill.transparency = 0.2
    circle1.line.fill.background()
    
    # Accent circles (top right)
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11), Inches(0.5),
        Inches(1.5), Inches(1.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = colors['accent']
    circle2.fill.transparency = 0.4
    circle2.line.fill.background()
    
    circle3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.8), Inches(1.5),
        Inches(1), Inches(1)
    )
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = colors['orange']
    circle3.fill.transparency = 0.3
    circle3.line.fill.background()
    
    # Modern accent bar (angled)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.5),
        Inches(0.2), Inches(2.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors['accent']
    accent_bar.line.fill.background()
    accent_bar.shadow.inherit = False
    
    # Title with shadow effect
    title_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.5),
        Inches(11), Inches(1.8)
    )
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = slide_data['title']
    
    p = text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.line_spacing = 1.15
    
    # Subtitle with highlight
    if slide_data['subtitle']:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(4.5),
            Inches(10), Inches(1)
        )
        text_frame = subtitle_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = slide_data['subtitle']
        
        p = text_frame.paragraphs[0]
        p.font.name = 'Segoe UI Light'
        p.font.size = Pt(30)
        p.font.color.rgb = colors['accent']
        p.line_spacing = 1.3
    
    # Key message in a styled box
    if slide_data['content'] and len(slide_data['content']) > 0:
        key_msg = slide_data['content'][0]
        if key_msg and not key_msg.startswith('-'):
            # Background box for key message
            msg_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.2), Inches(5.8),
                Inches(7), Inches(0.8)
            )
            msg_bg.fill.solid()
            msg_bg.fill.fore_color.rgb = colors['white']
            msg_bg.fill.transparency = 0.15
            msg_bg.line.fill.background()
            
            msg_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5.9),
                Inches(6.5), Inches(0.6)
            )
            text_frame = msg_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = key_msg
            
            p = text_frame.paragraphs[0]
            p.font.name = 'Segoe UI'
            p.font.size = Pt(22)
            p.font.color.rgb = colors['white']
            p.line_spacing = 1.3


def create_section_header_slide(prs, slide_data, slide_num, colors):
    """Create bold section header slide with dynamic shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Vibrant gradient background
    bg1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = colors['cyan']
    bg1.line.fill.background()
    
    # Overlapping gradient
    bg2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(0),
        Inches(8.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = colors['purple']
    bg2.fill.transparency = 0.5
    bg2.line.fill.background()
    
    # Large decorative circles
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9), Inches(3),
        Inches(5), Inches(5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = colors['accent']
    circle1.fill.transparency = 0.2
    circle1.line.fill.background()
    
    # Section number in a bold circle
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.5), Inches(2),
        Inches(2), Inches(2)
    )
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = colors['white']
    num_circle.fill.transparency = 0.15
    num_circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2),
        Inches(2), Inches(2)
    )
    text_frame = num_box.text_frame
    text_frame.text = f"{slide_num:02d}"
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    p.font.size = Pt(80)
    p.font.color.rgb = colors['white']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(4), Inches(2.5),
        Inches(8), Inches(2.5)
    )
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = slide_data['title']
    
    p = text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.line_spacing = 1.2
    
    # Accent bar under title
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4.8),
        Inches(4), Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = colors['accent']
    accent_bar.line.fill.background()
    
    add_footer(slide, slide_num, colors)


def create_content_slide(prs, slide_data, slide_num, colors):
    """Create visually engaging content slide with modern elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Subtle gradient background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['bg_light']
    background.line.fill.background()
    
    # Decorative side accent (left edge)
    side_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(7.5)
    )
    side_accent.fill.solid()
    side_accent.fill.fore_color.rgb = colors['blue']
    side_accent.line.fill.background()
    
    # Gradient header bar (two-tone effect)
    header_bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.15), Inches(0),
        Inches(8), Inches(1.1)
    )
    header_bar1.fill.solid()
    header_bar1.fill.fore_color.rgb = colors['blue']
    header_bar1.line.fill.background()
    
    header_bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.15), Inches(0),
        Inches(5.183), Inches(1.1)
    )
    header_bar2.fill.solid()
    header_bar2.fill.fore_color.rgb = colors['purple']
    header_bar2.line.fill.background()
    
    # Accent wave/strip
    accent_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.15), Inches(1.1),
        Inches(13.183), Inches(0.08)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = colors['accent']
    accent_strip.line.fill.background()
    
    # Decorative circles (top right corner)
    circle_deco = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.5), Inches(0.15),
        Inches(0.6), Inches(0.6)
    )
    circle_deco.fill.solid()
    circle_deco.fill.fore_color.rgb = colors['cyan']
    circle_deco.fill.transparency = 0.3
    circle_deco.line.fill.background()
    
    # Title in header with icon/number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.35), Inches(0.25),
        Inches(0.6), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = colors['accent']
    badge.line.fill.background()
    
    badge_text = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.25),
        Inches(0.6), Inches(0.6)
    )
    tf = badge_text.text_frame
    tf.text = str(slide_num)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['navy']
    
    # Title in header
    title_box = slide.shapes.add_textbox(
        Inches(1.1), Inches(0.25),
        Inches(11.5), Inches(0.6)
    )
    text_frame = title_box.text_frame
    text_frame.text = slide_data['title']
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # Content area with styled boxes for key items
    content_top = 1.6
    content_left = 0.6
    content_width = 12.2
    line_height = 0.5
    
    current_y = content_top
    item_count = 0
    
    # Define vibrant bullet colors
    bullet_colors = [colors['blue'], colors['purple'], colors['teal'], colors['orange']]
    
    for item in slide_data['content']:
        if not item or current_y > 6.5:  # Stop before footer
            break
        
        # Detect item type
        is_bullet = item.startswith('-') or item.startswith('•')
        is_numbered = re.match(r'^\d+\.', item)
        is_sub_bullet = item.startswith('  ')
        
        # Clean text
        clean_text = item.lstrip('-•').lstrip()
        clean_text = re.sub(r'^\d+\.\s*', '', clean_text)
        clean_text = clean_text.strip()
        
        if not clean_text:
            continue
        
        # Determine styling
        if is_sub_bullet:
            left_margin = content_left + 0.8
            font_size = 18
            color = colors['text_gray']
            bullet_shape = None
        elif is_bullet or is_numbered:
            left_margin = content_left + 0.5
            font_size = 21
            color = colors['text_dark']
            
            # Add colorful bullet circle
            bullet_color = bullet_colors[item_count % len(bullet_colors)]
            bullet_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(content_left + 0.2), Inches(current_y + 0.1),
                Inches(0.18), Inches(0.18)
            )
            bullet_circle.fill.solid()
            bullet_circle.fill.fore_color.rgb = bullet_color
            bullet_circle.line.fill.background()
            item_count += 1
        else:
            # Section header with accent bar
            left_margin = content_left + 0.2
            font_size = 24
            color = colors['blue']
            
            # Accent line before header
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(content_left + 0.2), Inches(current_y + 0.15),
                Inches(0.1), Inches(0.25)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = colors['accent']
            accent_line.line.fill.background()
            left_margin = content_left + 0.4
        
        # Create text box
        text_box = slide.shapes.add_textbox(
            Inches(left_margin), Inches(current_y),
            Inches(content_width - left_margin + content_left), Inches(line_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = clean_text
        
        p = text_frame.paragraphs[0]
        p.font.name = 'Segoe UI'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = 1.25
        
        # Bold for headers
        if not is_bullet and not is_numbered and not is_sub_bullet:
            p.font.bold = True
        
        current_y += line_height
    
    add_footer(slide, slide_num, colors)


def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_slides_professional.py input.md [output.pptx]")
        print("\nExample:")
        print("  python generate_slides_professional.py PFT_RFT_PM_Agent_Slides.md Professional.pptx")
        sys.exit(1)
    
    input_file = Path(sys.argv[1])
    if not input_file.exists():
        print(f"ERROR: Input file not found: {input_file}")
        sys.exit(1)
    
    output_file = Path(sys.argv[2]) if len(sys.argv) >= 3 else input_file.with_name('Professional_Presentation.pptx')
    
    print(f"Reading slides from: {input_file}")
    
    # Parse slides
    slides = parse_markdown_slides(input_file)
    print(f"✓ Parsed {len(slides)} slides")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    colors = set_professional_colors()
    
    slide_num = 0
    for idx, slide_data in enumerate(slides):
        slide_num = idx + 1
        
        print(f"  ✓ Creating slide {slide_num}: {slide_data['title'][:50]}...")
        
        if slide_data['type'] == 'title':
            create_title_slide(prs, slide_data, colors)
        elif 'Architecture' in slide_data['title'] or 'Roadmap' in slide_data['title']:
            create_section_header_slide(prs, slide_data, slide_num, colors)
        else:
            create_content_slide(prs, slide_data, slide_num, colors)
    
    # Save
    prs.save(str(output_file))
    
    file_size = output_file.stat().st_size / 1024
    print(f"\n✓ SUCCESS! Presentation created: {output_file}")
    print(f"  Slides: {len(slides)}")
    print(f"  Size: {file_size:.1f} KB")
    print(f"  Format: 16:9 widescreen")
    print(f"  Style: Professional modern design")
    
    # Auto-convert to PDF for VS Code viewing
    print(f"\n📄 Creating PDF for VS Code preview...")
    pdf_file = output_file.with_suffix('.pdf')
    try:
        import subprocess
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', str(output_file.parent), str(output_file)],
            capture_output=True,
            timeout=30
        )
        if pdf_file.exists():
            pdf_size = pdf_file.stat().st_size / 1024
            print(f"✓ PDF created: {pdf_file}")
            print(f"  Size: {pdf_size:.1f} KB")
            print(f"\n💡 Open in VS Code: code {pdf_file}")
        else:
            print(f"⚠ PDF conversion skipped (LibreOffice not available)")
    except Exception as e:
        print(f"⚠ PDF conversion skipped: {e}")


if __name__ == '__main__':
    main()
