#!/usr/bin/env python3
"""
Read and extract content from PowerPoint files (both .ppt and .pptx).
Extracts slide titles, text content, notes, and structure.
"""

from pptx import Presentation
from pathlib import Path
import json
import subprocess
import sys


def convert_ppt_to_pptx(ppt_path):
    """
    Convert old .ppt format to .pptx using LibreOffice.
    Falls back to extracting text with antiword or strings if conversion fails.
    
    Args:
        ppt_path: Path to the .ppt file
        
    Returns:
        Path to converted .pptx file or None if conversion failed
    """
    output_dir = ppt_path.parent
    
    # Try LibreOffice conversion
    try:
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'pptx',
            '--outdir', str(output_dir),
            str(ppt_path)
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        
        # Check if output file was created
        converted_file = output_dir / f"{ppt_path.stem}.pptx"
        if converted_file.exists() and converted_file != ppt_path:
            print(f"✅ Converted to: {converted_file}")
            return converted_file
    except (subprocess.TimeoutExpired, FileNotFoundError) as e:
        print(f"⚠️  LibreOffice conversion failed: {e}")
    
    return None


def extract_text_fallback(file_path):
    """
    Fallback method to extract text from old PowerPoint files.
    Uses strings command to extract readable text.
    """
    try:
        result = subprocess.run(['strings', str(file_path)], 
                              capture_output=True, text=True, timeout=10)
        text = result.stdout
        
        # Filter to likely content (longer strings, skip binary noise)
        lines = [line.strip() for line in text.split('\n') 
                if len(line.strip()) > 20 and not line.startswith('_')]
        
        return {
            "file_path": str(file_path),
            "method": "strings_extraction",
            "total_slides": "unknown",
            "extracted_text": lines[:100],  # First 100 meaningful lines
            "note": "This is a legacy .ppt file. Text extracted but structure lost."
        }
    except Exception as e:
        return {
            "file_path": str(file_path),
            "error": f"Failed to extract: {e}",
            "note": "File may be corrupted or inaccessible"
        }


def read_pptx(file_path):
    """
    Read a PPTX file and extract all content.
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        dict: Structured data with slides, titles, text, notes
    """
    prs = Presentation(file_path)
    
    data = {
        "file_path": str(file_path),
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "notes": "",
            "shapes": []
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            shape_info = {
                "type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                "text": ""
            }
            
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    shape_info["text"] = text
                    slide_data["content"].append(text)
                    
                    # Detect title (usually first text box or title placeholder)
                    if not slide_data["title"] and shape.shape_type == 1:  # Title placeholder
                        slide_data["title"] = text
            
            if shape_info["text"]:
                slide_data["shapes"].append(shape_info)
        
        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_data["notes"] = notes_slide.notes_text_frame.text.strip()
        
        # If no title was found, use first content as title
        if not slide_data["title"] and slide_data["content"]:
            slide_data["title"] = slide_data["content"][0]
        
        data["slides"].append(slide_data)
    
    return data


def print_pptx_content(data):
    """Print PPTX content in a readable format."""
    print("=" * 80)
    print(f"PowerPoint File: {data['file_path']}")
    
    if 'method' in data:
        print(f"Extraction Method: {data['method']}")
        print(f"Note: {data.get('note', '')}")
    
    if 'total_slides' in data:
        print(f"Total Slides: {data['total_slides']}")
    
    print("=" * 80)
    
    # Handle fallback text extraction
    if 'extracted_text' in data:
        print("\nExtracted Text (first 50 lines):")
        for i, line in enumerate(data['extracted_text'][:50], 1):
            print(f"{i}. {line}")
        return
    
    # Handle structured slide data
    if 'slides' in data:
        for slide in data["slides"]:
            print(f"\n{'='*80}")
            print(f"SLIDE {slide['slide_number']}: {slide['title']}")
            print(f"{'='*80}")
            
            if slide['content']:
                print("\nContent:")
                for i, text in enumerate(slide['content'], 1):
                    print(f"{i}. {text}")
            
            if slide['notes']:
                print(f"\nNotes:")
                print(slide['notes'])
            
            print()


def save_to_json(data, output_file):
    """Save extracted data to JSON file."""
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print(f"\n✅ Saved to {output_file}")


if __name__ == "__main__":
    # Read the CXA Workback Template
    pptx_path = Path("/Users/cyl/projects/scenara/src/workback_planning/examples/CXA/CXA Workback Template.pptx")
    
    if not pptx_path.exists():
        print(f"❌ File not found: {pptx_path}")
        exit(1)
    
    print(f"Reading: {pptx_path}")
    
    # Check if it's an old .ppt file (Composite Document)
    result = subprocess.run(['file', str(pptx_path)], capture_output=True, text=True)
    if 'Composite Document' in result.stdout:
        print("⚠️  Detected legacy .ppt format")
        
        # Try to convert
        converted = convert_ppt_to_pptx(pptx_path)
        if converted:
            pptx_path = converted
        else:
            print("❌ Conversion failed. Using text extraction fallback...")
            data = extract_text_fallback(pptx_path)
            print_pptx_content(data)
            
            output_json = pptx_path.parent / f"{pptx_path.stem}_extracted.json"
            save_to_json(data, output_json)
            exit(0)
    
    print()
    
    # Extract content from valid PPTX
    try:
        data = read_pptx(pptx_path)
        
        # Print to console
        print_pptx_content(data)
        
        # Save to JSON
        output_json = pptx_path.parent / f"{pptx_path.stem}_content.json"
        save_to_json(data, output_json)
        
        # Print summary
        print("\n" + "=" * 80)
        print("SUMMARY")
        print("=" * 80)
        print(f"Total slides: {data['total_slides']}")
        print(f"Slides with titles: {sum(1 for s in data['slides'] if s['title'])}")
        print(f"Slides with notes: {sum(1 for s in data['slides'] if s['notes'])}")
        print(f"Total text blocks: {sum(len(s['content']) for s in data['slides'])}")
    except Exception as e:
        print(f"❌ Error reading PowerPoint: {e}")
        print("\nTrying text extraction fallback...")
        data = extract_text_fallback(pptx_path)
        print_pptx_content(data)
        
        output_json = pptx_path.parent / f"{pptx_path.stem}_extracted.json"
        save_to_json(data, output_json)
